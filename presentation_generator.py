"""
presentation_generator.py
==========================
Generates a 5-7 slide executive PowerPoint presentation from monthly reports.
Uses python-pptx for slides and Gemini for content synthesis.
"""

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import GEMINI_API_KEY, GEMINI_MODEL, PRESENTATIONS


client = genai.Client(api_key=GEMINI_API_KEY)

# ── Brand Colours ──────────────────────────────────────────────────────────────
C_BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
C_BG_LIGHT   = RGBColor(0xF4, 0xF6, 0xF9)   # Off-white
C_ACCENT     = RGBColor(0x00, 0x8B, 0xFF)   # Bright blue
C_RED        = RGBColor(0xE5, 0x3E, 0x3E)
C_AMBER      = RGBColor(0xFF, 0xA5, 0x00)
C_GREEN      = RGBColor(0x28, 0xA7, 0x45)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0x8E, 0x99, 0xAB)
C_GOLD       = RGBColor(0xFF, 0xD7, 0x00)

RAG_COLORS   = {"Red": C_RED, "Amber": C_AMBER, "Green": C_GREEN}

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _add_bullet_textbox(slide, items: list[str], left, top, width, height,
                        font_size=11, color=C_WHITE, bullet_char="▸"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def _rag_badge(slide, rag: str, left: float, top: float):
    color = RAG_COLORS.get(rag, C_LIGHT_GRAY)
    badge = _add_rect(slide, left, top, 1.0, 0.35, color)
    _add_textbox(slide, rag.upper(), left + 0.05, top + 0.02,
                 0.9, 0.32, font_size=10, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)


def _set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Gemini synthesis
# ─────────────────────────────────────────────────────────────────────────────

SYNTHESIS_PROMPT = """\
You are a Principal Consultant preparing an executive briefing for a VP to present
to a client. Analyse these project health reports and produce a monthly synthesis.

=== REPORTS ===
{reports_json}

Today: {today}

Return ONLY a valid JSON object with this schema:
{{
  "month_label": "July 2026",
  "portfolio_rag": "Red|Amber|Green",
  "headline": "One punchy headline sentence for the VP",
  "portfolio_snapshot": "2-3 sentence portfolio overview with specific numbers",
  "trend_analysis": [
    "Trend insight 1 (with evidence)",
    "Trend insight 2",
    "Trend insight 3"
  ],
  "emerging_risks": [
    {{"risk": "Risk title", "detail": "1-sentence detail", "severity": "High|Medium"}},
    {{"risk": "...", "detail": "...", "severity": "..."}}
  ],
  "project_summaries": [
    {{"name": "...", "rag": "Red|Amber|Green", "one_liner": "...", "action": "..."}},
    ...
  ],
  "executive_recommendations": [
    "Recommendation 1 (specific and actionable)",
    "Recommendation 2",
    "Recommendation 3"
  ],
  "next_steps": [
    "Next step 1",
    "Next step 2",
    "Next step 3"
  ]
}}
"""


def synthesise_with_gemini(reports: list[dict]) -> dict:
    """Call Gemini to synthesise multiple reports into executive content."""
    # Strip heavy meta to save tokens
    slim_reports = []
    for r in reports:
        slim_reports.append({
            "project_name":      r.get("project_name"),
            "report_date":       r.get("report_date"),
            "rag_status":        r.get("rag_status"),
            "executive_summary": r.get("executive_summary"),
            "reasoning":         r.get("reasoning", {}),
            "key_metrics":       r.get("key_metrics", {}),
            "top_risks":         r.get("top_risks", []),
            "recommendations":   r.get("recommendations", []),
        })

    today  = datetime.now().strftime("%d %B %Y")
    prompt = SYNTHESIS_PROMPT.format(
        reports_json=json.dumps(slim_reports, indent=2, default=str),
        today=today,
    )

    response = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=prompt,
    )
    raw      = response.text.strip()
    raw      = re.sub(r"^```[a-z]*\s*", "", raw, flags=re.IGNORECASE)
    raw      = re.sub(r"\s*```$", "", raw)
    return json.loads(raw)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def _slide_title(prs: Presentation, synthesis: dict):
    """Slide 1 – Title / Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, C_BG_DARK)

    # Top accent bar
    _add_rect(slide, 0, 0, 13.33, 0.12, C_ACCENT)

    # Logo placeholder / company name
    _add_textbox(slide, "ZYCUS", 0.4, 0.25, 3, 0.5,
                 font_size=13, bold=True, color=C_ACCENT)

    # Main title
    _add_textbox(slide, "Project Portfolio\nHealth Report",
                 1.2, 1.8, 11, 1.8, font_size=44, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

    # Headline
    _add_textbox(slide, synthesis.get("headline", ""),
                 1.2, 3.7, 11, 0.8, font_size=16, italic=True,
                 color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Month & RAG badge area
    month = synthesis.get("month_label", datetime.now().strftime("%B %Y"))
    _add_textbox(slide, month, 0.4, 6.5, 5, 0.5,
                 font_size=14, color=C_LIGHT_GRAY)

    portfolio_rag = synthesis.get("portfolio_rag", "Amber")
    rag_color = RAG_COLORS.get(portfolio_rag, C_AMBER)
    _add_rect(slide, 10.5, 6.3, 2.5, 0.7, rag_color)
    _add_textbox(slide, f"Portfolio Status: {portfolio_rag}",
                 10.5, 6.35, 2.5, 0.6, font_size=13, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

    # Bottom bar
    _add_rect(slide, 0, 7.38, 13.33, 0.12, C_ACCENT)


def _slide_portfolio_snapshot(prs: Presentation, synthesis: dict, reports: list[dict]):
    """Slide 2 – Portfolio Snapshot with RAG tiles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_DARK)
    _add_rect(slide, 0, 0, 13.33, 0.12, C_ACCENT)

    # Title
    _add_textbox(slide, "Portfolio Snapshot", 0.4, 0.2, 10, 0.6,
                 font_size=24, bold=True, color=C_WHITE)
    _add_textbox(slide, synthesis.get("month_label", ""),
                 10.5, 0.25, 2.5, 0.5, font_size=13,
                 color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # Overview text
    snapshot = synthesis.get("portfolio_snapshot", "")
    _add_textbox(slide, snapshot, 0.4, 0.95, 12.5, 1.0,
                 font_size=13, color=C_LIGHT_GRAY)

    # Project RAG tiles
    projects = synthesis.get("project_summaries", [])
    tile_w, tile_h = 3.8, 2.5
    cols = min(len(projects), 3)
    start_x = (13.33 - cols * tile_w - (cols - 1) * 0.25) / 2

    for i, proj in enumerate(projects[:4]):
        col = i % 3
        row = i // 3
        x = start_x + col * (tile_w + 0.25)
        y = 2.1 + row * (tile_h + 0.2)
        rag = proj.get("rag", "Amber")
        rag_color = RAG_COLORS.get(rag, C_AMBER)

        # Tile background
        _add_rect(slide, x, y, tile_w, tile_h, RGBColor(0x14, 0x26, 0x3A))

        # Top colour bar (RAG indicator)
        _add_rect(slide, x, y, tile_w, 0.12, rag_color)

        # RAG label
        _add_textbox(slide, rag.upper(), x + 0.1, y + 0.18,
                     0.8, 0.3, font_size=9, bold=True, color=rag_color)

        # Project name
        _add_textbox(slide, proj.get("name", "")[:35], x + 0.15, y + 0.5,
                     tile_w - 0.3, 0.45, font_size=12, bold=True, color=C_WHITE)

        # One-liner
        _add_textbox(slide, proj.get("one_liner", "")[:100],
                     x + 0.15, y + 1.0, tile_w - 0.3, 0.8,
                     font_size=10, color=C_LIGHT_GRAY)

        # Action
        _add_textbox(slide, f"▶ {proj.get('action', '')[:80]}",
                     x + 0.15, y + 1.85, tile_w - 0.3, 0.5,
                     font_size=9, color=C_ACCENT, italic=True)

    _add_rect(slide, 0, 7.38, 13.33, 0.12, C_ACCENT)


def _slide_trends(prs: Presentation, synthesis: dict):
    """Slide 3 – Trend Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_DARK)
    _add_rect(slide, 0, 0, 13.33, 0.12, C_ACCENT)

    _add_textbox(slide, "Trend Analysis", 0.4, 0.2, 10, 0.6,
                 font_size=24, bold=True, color=C_WHITE)

    trends = synthesis.get("trend_analysis", [])

    icons = ["[UP]", "[DOWN]", "️", "", "[CYCLE]"]
    for i, trend in enumerate(trends[:5]):
        y = 1.0 + i * 1.1
        # Number circle
        _add_rect(slide, 0.4, y, 0.55, 0.55, C_ACCENT)
        _add_textbox(slide, str(i + 1), 0.4, y + 0.06, 0.55, 0.44,
                     font_size=14, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER)
        # Trend text
        _add_textbox(slide, trend, 1.1, y, 11.5, 0.9,
                     font_size=13, color=C_WHITE)

    _add_rect(slide, 0, 7.38, 13.33, 0.12, C_ACCENT)


def _slide_risks(prs: Presentation, synthesis: dict):
    """Slide 4 – Emerging Risks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_DARK)
    _add_rect(slide, 0, 0, 13.33, 0.12, C_ACCENT)

    _add_textbox(slide, "Emerging Risks", 0.4, 0.2, 10, 0.6,
                 font_size=24, bold=True, color=C_WHITE)

    risks = synthesis.get("emerging_risks", [])
    sev_colors = {"High": C_RED, "Medium": C_AMBER, "Low": C_GREEN}

    for i, risk in enumerate(risks[:5]):
        y = 1.0 + i * 1.1
        sev   = risk.get("severity", "Medium")
        color = sev_colors.get(sev, C_AMBER)

        # Severity badge
        _add_rect(slide, 0.4, y + 0.08, 1.1, 0.35, color)
        _add_textbox(slide, sev.upper(), 0.4, y + 0.1, 1.1, 0.3,
                     font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Risk title
        _add_textbox(slide, risk.get("risk", ""), 1.65, y,
                     5.5, 0.45, font_size=13, bold=True, color=C_WHITE)

        # Detail
        _add_textbox(slide, risk.get("detail", ""), 1.65, y + 0.45,
                     11.0, 0.5, font_size=11, color=C_LIGHT_GRAY)

    _add_rect(slide, 0, 7.38, 13.33, 0.12, C_ACCENT)


def _slide_recommendations(prs: Presentation, synthesis: dict):
    """Slide 5 – Executive Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_DARK)
    _add_rect(slide, 0, 0, 13.33, 0.12, C_ACCENT)

    _add_textbox(slide, "Executive Recommendations", 0.4, 0.2, 12, 0.6,
                 font_size=24, bold=True, color=C_WHITE)

    recs = synthesis.get("executive_recommendations", [])
    label_colors = [C_RED, C_AMBER, C_ACCENT, C_GREEN, C_GREEN]

    for i, rec in enumerate(recs[:5]):
        y = 1.0 + i * 1.1
        _add_rect(slide, 0.4, y, 0.55, 0.55, label_colors[i])
        _add_textbox(slide, str(i + 1), 0.4, y + 0.06, 0.55, 0.44,
                     font_size=14, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide, rec, 1.1, y, 11.8, 0.9,
                     font_size=13, color=C_WHITE)

    _add_rect(slide, 0, 7.38, 13.33, 0.12, C_ACCENT)


def _slide_next_steps(prs: Presentation, synthesis: dict):
    """Slide 6 – Next Steps & Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_DARK)
    _add_rect(slide, 0, 0, 13.33, 0.12, C_ACCENT)

    _add_textbox(slide, "Next Steps", 0.4, 0.2, 10, 0.6,
                 font_size=24, bold=True, color=C_WHITE)

    # Two-column layout: Left = next steps, Right = closing note
    next_steps = synthesis.get("next_steps", [])
    _add_bullet_textbox(slide, next_steps, 0.4, 1.0, 7.5, 5.5,
                        font_size=14, color=C_WHITE)

    # Right panel – portfolio RAG summary
    _add_rect(slide, 8.8, 1.0, 4.2, 2.2, RGBColor(0x14, 0x26, 0x3A))
    port_rag = synthesis.get("portfolio_rag", "Amber")
    rag_color = RAG_COLORS.get(port_rag, C_AMBER)
    _add_rect(slide, 8.8, 1.0, 4.2, 0.12, rag_color)
    _add_textbox(slide, "Portfolio Status", 8.9, 1.15, 4.0, 0.45,
                 font_size=11, color=C_LIGHT_GRAY, bold=True)
    _add_textbox(slide, port_rag, 8.9, 1.6, 4.0, 0.7,
                 font_size=32, bold=True, color=rag_color,
                 align=PP_ALIGN.CENTER)

    # Confidentiality footer
    _add_textbox(slide, "CONFIDENTIAL – For executive use only",
                 0.4, 7.0, 12.5, 0.3, font_size=9,
                 color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    _add_rect(slide, 0, 7.38, 13.33, 0.12, C_ACCENT)


# ─────────────────────────────────────────────────────────────────────────────
# Main entry point
# ─────────────────────────────────────────────────────────────────────────────

def generate_presentation(reports: list[dict],
                          output_path: Optional[Path] = None) -> Path:
    """
    Generate a 6-slide executive PowerPoint from a list of project reports.

    Parameters
    ----------
    reports     : list of report dicts (from gemini_agent.analyse_project)
    output_path : optional custom save path

    Returns
    -------
    Path to saved .pptx file
    """
    if not reports:
        raise ValueError("No reports provided for presentation generation.")

    # 1. Synthesise with Gemini
    print("[AI] Synthesising monthly insights with Gemini...")
    synthesis = synthesise_with_gemini(reports)

    # 2. Build presentation
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("[CHART] Building slides...")
    _slide_title(prs, synthesis)
    _slide_portfolio_snapshot(prs, synthesis, reports)
    _slide_trends(prs, synthesis)
    _slide_risks(prs, synthesis)
    _slide_recommendations(prs, synthesis)
    _slide_next_steps(prs, synthesis)

    # 3. Save
    if output_path is None:
        month = datetime.now().strftime("%Y%m")
        output_path = PRESENTATIONS / f"Portfolio_Health_Report_{month}.pptx"

    prs.save(str(output_path))
    print(f" Presentation saved: {output_path}")
    return output_path
